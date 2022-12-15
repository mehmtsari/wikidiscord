from random import randint
from os import getcwd,remove
from time import sleep
import wikipedia



path = str(getcwd())
file_path = path+"/"+"content/"
file_path = str(file_path)
required_path = path+"/"+"required/"
required_path = str(required_path)


def search_connect(search):
    try:
        wikipedia.set_lang("tr")
        voice_content = wikipedia.summary(search)
        return True
    except:
        return False
    
    
def information_voice(search):
    from gtts import gTTS


    #Ses texti
    try:
        wikipedia.set_lang("tr")
        voice_content = wikipedia.summary(search)
        voice_text = "Bu Bilgilendirme Mesajı Wikipedia Asistan Tarafından Oluşturulmuştur...\n"
        voice_text += str(voice_content)
    except:
        print("System: {} İsimli İçerik Bulunamadı.".format(search))
        sleep(0.4)
        return
    #Ses Kaydetme
    try:
        rand = randint(1,10000)
        rand = str(rand)
        voice = gTTS(text=voice_text,lang="tr",slow=False)
        voicename=""
        voicename = "voice-"+str(rand)+".mp4"
        voice_path = path+"/"+"content"
        voice.save(file_path+voicename)
        voice_x_path = file_path+voicename
        return voice_x_path
    except:
        print("System error: Ses Dosyası Oluşturulamadı.")
        



def main_wiki_func(search):
    contentpath = path+"\content"

    rand = randint(1,10000)
    rand = str(rand)
    
    #SES Dosyası---------------------------------------------------------------------------------------------------------------------------
    from gtts import gTTS


    #Ses texti
    try:
        wikipedia.set_lang("tr")
        voice_content = wikipedia.summary(search)
        voice_text = "Bu Bilgilendirme Mesajı Wikipedia Asistan Tarafından Oluşturulmuştur...\n"
        voice_text += str(voice_content)
    except:
        print("System: {} İsimli İçerik Bulunamadı.".format(search))
        sleep(0.4)
        return
    #Ses Kaydetme
    try:
        voice = gTTS(text=voice_text,lang="tr",slow=False)
        voicename=""
        voicename = "voice-"+str(rand)+".mp3"
        voice_path = path+"/"+"content"
        voice.save(file_path+voicename)
        voice_x_path = file_path+voicename
    except:
        print("System error: Ses Dosyası Oluşturulamadı.")
    #Bütün bir wikipedia sayfası---------------------------------------------------------------------------------------------------------------------
    try:
        wikipedia.set_lang("tr")
        page = wikipedia.page(search)
        maintitle = page.title
        content = page.content

        content = str(content)
        contentlist= []
        contentlist = content.split("==")
        contentlist.insert(0,maintitle)
        title_list = []
        content_list=[]

        x = 0
        for i in contentlist:
            changed = i.strip("=")
            changed = changed.strip("")
            changed = changed.strip(" ")
            if x % 2 == 0:
            
                title_list.append(i)
            else:
                content_list.append(i)
            x += 1
    except:
        print("System: {} İsimli İçerik Bulunamadı.".format(search))
        sleep(0.4)
        return
        

    #SLAYT---------------------------------------------------------------------------------------------------------------------------------
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    main_slayt_path = required_path +"/main/Main_Slayt001(Dont Delete).pptx"
    prs = Presentation(main_slayt_path)
    #BAŞLIK
    slide_entry1 = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_entry1)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = maintitle
    subtitle1.text = "MEHMET SARI"



    #İÇERİK VE BAŞLIKLARI

    for i,j in zip(title_list,content_list):
        
        # İÇERİK BAŞLIĞI
        i = i.replace("=","")
        i = i.replace("==","")
        i = i.strip()
        
        # İÇERİK
        old_contents= []
        old_contents = j.split("\n")
        
        contents = []
        index = 0
        for k in old_contents:
            
            lenght = (len(k))
            lenght = int(lenght)
            if lenght < 20 :
                pass
            else:
                contents.append((k))
        
        old_contents.clear()

        
        # HER BİR PARAGRAF
        for p in contents:
            p = str(p)
            plenght = (len(p))
            plenght = int(plenght)
            
            one = int(0)
            two = int(one +1)
            three = int(one +2)
        
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            placeholder_title = slide.placeholders[0]
            placeholder_title.text = i
            s = placeholder_title.text_frame.paragraphs[0]
            s.font.size = Pt(56)
            s.font.bold = True
            
            placeholder = slide.placeholders[1]
            try:
                tf = placeholder.text_frame
                t = tf.add_paragraph()
                t.text = p
                t.font.size = Pt(18)
            except:
                continue
    t = tf.add_paragraph()
    t.text = "Bu yazı Wikipedia Asistan Tarafından Oluşturulmuştur."
    t.font.size = Pt(40)
    pptname = ("slayt-"+rand+".pptx")
    prs.save(file_path+pptname)
    pptx_x_path = file_path+pptname

    #docx Dosyası---------------------------------------------------------------------------------------------------------------------------------

    try:
        docxname = ("word-"+rand+".docx")
        from docx import Document
        document = Document()

        document.add_heading(maintitle, 0)
        for i,j in zip(title_list,content_list):
            
            # İÇERİK BAŞLIĞI
            i = i.replace("=","")
            i = i.replace("==","")
            i = i.strip()
            
            # İÇERİK
            old_contents= []
            old_contents = j.split("\n")
            
            contents = []
            index = 0
            for k in old_contents:
                
                lenght = (len(k))
                lenght = int(lenght)
                if lenght < 20 :
                    pass
                else:
                    contents.append((k))
            
            old_contents.clear()

            
            # HER BİR PARAGRAF
            document.add_paragraph(i)
            pgr = ""
            for p in contents:
                text = p+"\n"
                pgr += text
            
            document.add_paragraph(str(pgr))
        document.add_paragraph("Bu yazı Wikipedia Asistan Tarafından Oluşturulmuştur.")
        document.save(file_path + docxname)
        
        docx_x_path = file_path + docxname
    except:
        print("Docx Dosyası Oluşturulamadı")
    
    return (voice_x_path+"#"+pptx_x_path+"#"+docx_x_path)
